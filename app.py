import os
import io
import re
import json
import math
import time
import tempfile
from collections import Counter, defaultdict

import cv2
import faiss
import librosa
import numpy as np
import streamlit as st
import shutil
import imageio_ffmpeg

# Setup ffmpeg so we don't need a system installation
ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
ffmpeg_alias_dir = tempfile.mkdtemp(prefix="ffmpeg_alias_")
ffmpeg_exe_name = "ffmpeg.exe" if os.name == 'nt' else "ffmpeg"
ffmpeg_symlink = os.path.join(ffmpeg_alias_dir, ffmpeg_exe_name)
if not os.path.exists(ffmpeg_symlink):
    shutil.copyfile(ffmpeg_path, ffmpeg_symlink)
    if os.name != 'nt':
        os.chmod(ffmpeg_symlink, 0o755)
os.environ["PATH"] = ffmpeg_alias_dir + os.pathsep + os.environ.get("PATH", "")
from deepface import DeepFace
from moviepy import VideoFileClip
from sentence_transformers import SentenceTransformer, util

# PPT / PDF parsing
from pptx import Presentation
import fitz  # PyMuPDF

# Whisper
import whisper
from google import genai
from pydantic import BaseModel, Field

# =========================================================
# CONFIG
# =========================================================
st.set_page_config(page_title="GenAI Viva Coach", layout="wide")

GEMINI_API_KEY_1 = "FIRST_API_KEY_HERE"
GEMINI_API_KEY_2 = "SECOND_API_KEY_HERE"

class VivaQuestion(BaseModel):
    type: str = Field(description="The category of the question (e.g., 'retrieval_gap', 'nervousness', 'content_depth', 'general').")
    question: str = Field(description="The dynamically generated viva question.")

class VivaQuestionList(BaseModel):
    questions: list[VivaQuestion]

class AnswerEvaluation(BaseModel):
    relevance: float = Field(description="Score out of 100 for correctness/relevance to the slide context.")
    depth: float = Field(description="Score out of 100 for proper technical depth.")
    clarity: float = Field(description="Score out of 100 for clarity and coherence.")
    score: float = Field(description="The overall unified score out of 100.")
    feedback: str = Field(description="Constructive and professional feedback on how this answer could be improved.")


EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
WHISPER_MODEL_NAME = "base"
FRAME_SAMPLE_SECONDS = 2
MAX_QUESTIONS = 6
FILLER_WORDS = [
    "um", "uh", "like", "basically", "actually", "you know", "sort of", "kind of"
]



# =========================================================
# CACHING
# =========================================================
@st.cache_resource
def load_embedder():
    return SentenceTransformer(EMBED_MODEL_NAME)


@st.cache_resource
def load_whisper_model():
    return whisper.load_model(WHISPER_MODEL_NAME)


# =========================================================
# FILE HELPERS
# =========================================================
def save_uploaded_file(uploaded_file, suffix):
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getbuffer())
        return tmp.name


# =========================================================
# PPT / PDF EXTRACTION
# =========================================================
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slide_chunks = []

    for idx, slide in enumerate(prs.slides, start=1):
        collected = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    txt = shape.text.strip()
                    if txt:
                        collected.append(txt)
                elif hasattr(shape, "table") and shape.table is not None:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            txt = cell.text.strip()
                            if txt:
                                collected.append(txt)
            except Exception:
                continue

        slide_text = "\n".join(collected).strip()
        if slide_text:
            slide_chunks.append({
                "source": f"Slide {idx}",
                "text": slide_text,
            })

    return slide_chunks



def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    page_chunks = []

    for i in range(len(doc)):
        page = doc[i]
        text = page.get_text("text").strip()
        if text:
            page_chunks.append({
                "source": f"Page {i + 1}",
                "text": text,
            })
    doc.close()
    return page_chunks


# =========================================================
# VIDEO + AUDIO
# =========================================================
def extract_audio_from_video(video_path, audio_out_path):
    import subprocess
    cmd = [
        "ffmpeg", 
        "-y", 
        "-i", video_path, 
        "-vn", 
        "-acodec", "pcm_s16le", 
        "-ar", "16000", 
        "-ac", "1", 
        audio_out_path
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)



def get_video_meta(video_path):
    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 25.0
    frame_count = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0
    duration = frame_count / fps if fps > 0 else 0
    width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
    height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
    cap.release()
    return {
        "fps": fps,
        "frame_count": int(frame_count),
        "duration": duration,
        "width": width,
        "height": height,
    }


# =========================================================
# TRANSCRIPTION
# =========================================================
def transcribe_audio(audio_path):
    model = load_whisper_model()
    result = model.transcribe(audio_path)
    text = result.get("text", "").strip()
    segments = result.get("segments", [])
    return text, segments


# =========================================================
# SPEECH ANALYSIS
# =========================================================
def analyze_speech(audio_path, transcript_text, segments, slide_text_full):
    y, sr = librosa.load(audio_path, sr=None)
    duration = librosa.get_duration(y=y, sr=sr)

    rms = librosa.feature.rms(y=y)[0]
    energy_mean = float(np.mean(rms)) if len(rms) else 0.0
    energy_std = float(np.std(rms)) if len(rms) else 0.0

    words = re.findall(r"\b\w+\b", transcript_text.lower())
    word_count = len(words)
    words_per_minute = (word_count / max(duration, 1e-6)) * 60.0

    filler_count = 0
    lowered = transcript_text.lower()
    for filler in FILLER_WORDS:
        filler_count += len(re.findall(rf"\b{re.escape(filler)}\b", lowered))

    filler_rate = filler_count / max(word_count, 1)

    long_pauses = 0
    if segments:
        prev_end = None
        for seg in segments:
            start = float(seg.get("start", 0.0))
            end = float(seg.get("end", 0.0))
            if prev_end is not None and (start - prev_end) > 1.5:
                long_pauses += 1
            prev_end = end

    # Content coverage: how much of slide keywords appeared in transcript
    slide_keywords = extract_keywords(slide_text_full)
    transcript_words = set(re.findall(r"\b[a-zA-Z][a-zA-Z0-9_-]*\b", transcript_text.lower()))
    overlap = len(slide_keywords & transcript_words)
    coverage = overlap / max(len(slide_keywords), 1)

    # Fluency score
    wpm_score = score_words_per_minute(words_per_minute)
    filler_score = max(0.0, 100.0 - filler_rate * 800.0)
    pause_score = max(0.0, 100.0 - long_pauses * 8.0)
    fluency_score = clamp((0.4 * wpm_score) + (0.35 * filler_score) + (0.25 * pause_score))

    # Vocal confidence score
    vocal_confidence = clamp((energy_mean * 2200.0) - (energy_std * 600.0) + 35.0)

    # Content delivery score
    content_delivery = clamp((coverage * 100.0 * 0.7) + min(word_count / 6.0, 30.0))

    return {
        "duration_seconds": duration,
        "word_count": word_count,
        "words_per_minute": words_per_minute,
        "filler_count": filler_count,
        "filler_rate": filler_rate,
        "long_pauses": long_pauses,
        "coverage": coverage,
        "energy_mean": energy_mean,
        "energy_std": energy_std,
        "fluency_score": fluency_score,
        "vocal_confidence": vocal_confidence,
        "content_delivery": content_delivery,
    }



def score_words_per_minute(wpm):
    # best around ~110 to 170 wpm
    ideal = 140.0
    diff = abs(wpm - ideal)
    score = 100.0 - min(diff * 1.2, 100.0)
    return clamp(score)



def extract_keywords(text, max_keywords=80):
    tokens = re.findall(r"[a-zA-Z][a-zA-Z0-9_-]*", text.lower())
    stop = {
        "i", "me", "my", "myself", "we", "our", "ours", "ourselves", "you", "your", "yours",
        "yourself", "yourselves", "he", "him", "his", "himself", "she", "her", "hers",
        "herself", "it", "its", "itself", "they", "them", "their", "theirs", "themselves",
        "what", "which", "who", "whom", "this", "that", "these", "those", "am", "is", "are",
        "was", "were", "be", "been", "being", "have", "has", "had", "having", "do", "does",
        "did", "doing", "a", "an", "the", "and", "but", "if", "or", "because", "as", "until",
        "while", "of", "at", "by", "for", "with", "about", "against", "between", "into",
        "through", "during", "before", "after", "above", "below", "to", "from", "up", "down",
        "in", "out", "on", "off", "over", "under", "again", "further", "then", "once", "here",
        "there", "when", "where", "why", "how", "all", "any", "both", "each", "few", "more",
        "most", "other", "some", "such", "no", "nor", "not", "only", "own", "same", "so",
        "than", "too", "very", "s", "t", "can", "will", "just", "don", "should", "now", "using", "used"
    }
    counts = Counter([t for t in tokens if t not in stop and len(t) > 2])
    top = [word for word, _ in counts.most_common(max_keywords)]
    return set(top)


# =========================================================
# EMOTION ANALYSIS
# =========================================================
def analyze_video_emotions(video_path):
    meta = get_video_meta(video_path)
    fps = meta["fps"] if meta["fps"] > 0 else 25.0
    frame_step = max(1, int(fps * FRAME_SAMPLE_SECONDS))

    cap = cv2.VideoCapture(video_path)
    frame_idx = 0
    analyzed = 0
    emotion_acc = defaultdict(float)
    no_face_frames = 0
    dominant_emotions = []

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break

        if frame_idx % frame_step == 0:
            try:
                result = DeepFace.analyze(frame, actions=["emotion"], enforce_detection=False)
                if isinstance(result, list):
                    result = result[0]
                emotions = result.get("emotion", {})
                if emotions:
                    analyzed += 1
                    for k, v in emotions.items():
                        emotion_acc[k.lower()] += float(v)
                    dominant_emotions.append(max(emotions, key=emotions.get).lower())
                else:
                    no_face_frames += 1
            except Exception:
                no_face_frames += 1

        frame_idx += 1

    cap.release()

    avg_emotions = {}
    if analyzed > 0:
        for k, v in emotion_acc.items():
            avg_emotions[k] = v / analyzed

    fear = avg_emotions.get("fear", 0.0)
    sad = avg_emotions.get("sad", 0.0)
    happy = avg_emotions.get("happy", 0.0)
    neutral = avg_emotions.get("neutral", 0.0)
    surprise = avg_emotions.get("surprise", 0.0)
    disgust = avg_emotions.get("disgust", 0.0)

    face_visibility_ratio = analyzed / max(analyzed + no_face_frames, 1)
    facial_confidence = clamp(
        (0.45 * neutral)
        + (0.25 * happy)
        + (face_visibility_ratio * 20.0)
        - (0.35 * fear)
        - (0.20 * sad)
        - (0.10 * surprise)
        - (0.10 * disgust)
        + 30.0
    )

    return {
        "video_meta": meta,
        "analyzed_frames": analyzed,
        "no_face_frames": no_face_frames,
        "face_visibility_ratio": face_visibility_ratio,
        "avg_emotions": avg_emotions,
        "dominant_emotions": Counter(dominant_emotions),
        "facial_confidence": facial_confidence,
    }


# =========================================================
# RETRIEVAL / RAG
# =========================================================
def chunk_transcript(segments):
    chunks = []
    if not segments:
        return chunks

    buffer = []
    start_time = None
    for seg in segments:
        text = seg.get("text", "").strip()
        if not text:
            continue
        if start_time is None:
            start_time = seg.get("start", 0.0)
        buffer.append(text)

        if len(" ".join(buffer).split()) >= 70:
            end_time = seg.get("end", 0.0)
            chunks.append({
                "source": f"Transcript {start_time:.1f}s-{end_time:.1f}s",
                "text": " ".join(buffer).strip(),
            })
            buffer = []
            start_time = None

    if buffer:
        chunks.append({
            "source": f"Transcript {start_time:.1f}s-end",
            "text": " ".join(buffer).strip(),
        })

    return chunks



def build_vector_index(chunks):
    embedder = load_embedder()
    texts = [c["text"] for c in chunks]
    embeddings = embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)
    index = faiss.IndexFlatIP(embeddings.shape[1])
    index.add(embeddings.astype(np.float32))
    return embedder, index, embeddings



def retrieve_context(query, embedder, index, chunks, k=3):
    q = embedder.encode([query], convert_to_numpy=True, normalize_embeddings=True)
    scores, ids = index.search(q.astype(np.float32), min(k, len(chunks)))
    results = []
    for score, idx in zip(scores[0], ids[0]):
        if idx == -1:
            continue
        results.append({
            "score": float(score),
            "source": chunks[idx]["source"],
            "text": chunks[idx]["text"],
        })
    return results


# =========================================================
# QUESTION GENERATION
# =========================================================
def detect_weak_areas(emotion_data, speech_data):
    weak_areas = []

    fear = emotion_data["avg_emotions"].get("fear", 0.0)
    face_visibility = emotion_data["face_visibility_ratio"]
    coverage = speech_data["coverage"]
    filler_rate = speech_data["filler_rate"]
    wpm = speech_data["words_per_minute"]
    content_delivery = speech_data["content_delivery"]

    if fear > 18 or face_visibility < 0.6:
        weak_areas.append("nervousness")

    if filler_rate > 0.02 or wpm < 90 or wpm > 185:
        weak_areas.append("fluency")

    if coverage < 0.35 or content_delivery < 45:
        weak_areas.append("content_depth")

    if coverage < 0.28:
        weak_areas.append("retrieval_gap")

    if speech_data["fluency_score"] < 55 or speech_data["content_delivery"] < 50:
        weak_areas.append("technical_understanding")

    # unique, preserve order
    unique = []
    seen = set()
    for area in weak_areas:
        if area not in seen:
            unique.append(area)
            seen.add(area)
    return unique



def generate_viva_questions(project_text, weak_areas, embedder=None, index=None, chunks=None, api_key=""):
    if len(project_text) > 8000:
        project_text = project_text[:8000] + "\n...[truncated for token limits]"
        
    if not api_key or api_key == "YOUR_API_KEY_HERE" or api_key == "YOUR_SECOND_API_KEY_HERE":
        st.error("Please add your Gemini API Key in app.py to generate dynamic questions.")
        st.stop()
        
    client = genai.Client(api_key=api_key)
    prompt = f"""
You are an expert capstone presentation evaluator.
Based on the following presentation transcript and slide text, generate up to {MAX_QUESTIONS} personalized viva questions.
The speaker's weak areas were identified as: {', '.join(weak_areas) if weak_areas else 'none major'}.
Focus questions on concepts that were under-explained in the transcript compared to the slides.

Transcript & Slides Content:
{project_text}
"""
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config={
                'response_mime_type': 'application/json',
                'response_schema': VivaQuestionList,
            },
        )
        parsed = json.loads(response.text)
        return [{"type": q["type"], "question": q["question"]} for q in parsed.get("questions", [])]
    except Exception as e:
        print(f"Error generating questions: {e}")
        return [{"type": "fallback", "question": "Could you elaborate on the core problem your project solves?"}]


# =========================================================
# ANSWER EVALUATION
# =========================================================
def evaluate_answer(answer, question_obj, embedder=None, index=None, chunks=None, slide_text="", api_key=""):
    answer = (answer or "").strip()
    if not answer:
        return {
            "relevance": 0.0,
            "depth": 0.0,
            "clarity": 0.0,
            "score": 0.0,
            "feedback": "No answer provided.",
        }

    if not api_key or api_key == "YOUR_API_KEY_HERE" or api_key == "YOUR_SECOND_API_KEY_HERE":
        return {"relevance": 0.0, "depth": 0.0, "clarity": 0.0, "score": 0.0, "feedback": "API Key missing."}

    client = genai.Client(api_key=api_key)
    question_text = question_obj.get("question", str(question_obj))
    
    context_text = slide_text
    if embedder and index and chunks:
        relevant_docs = retrieve_context(question_text, embedder, index, chunks, k=4)
        context_text = "\n".join([doc["text"] for doc in relevant_docs])

    prompt = f"""
You are grading a student's answer to a viva question based on their project slides.
Question: {question_text}
Student's Answer: {answer}
Relevant Project Context: {context_text}

Provide an honest score out of 100 for correctness/relevance, technical depth, and clarity. Compute an overall unified score. Give constructive feedback.
"""
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config={
                'response_mime_type': 'application/json',
                'response_schema': AnswerEvaluation,
            },
        )
        data = json.loads(response.text)
        return {
            "relevance": data.get("relevance", 0.0),
            "depth": data.get("depth", 0.0),
            "clarity": data.get("clarity", 0.0),
            "score": data.get("score", 0.0),
            "feedback": data.get("feedback", "Could not parse feedback."),
        }
    except Exception as e:
        print(f"Error evaluating answer: {e}")
        return {
            "relevance": 0.0, "depth": 0.0, "clarity": 0.0, "score": 0.0, "feedback": "Evaluation failed due to an error."
        }


# =========================================================
# SCORING
# =========================================================
def clamp(x, lo=0.0, hi=100.0):
    return max(lo, min(hi, float(x)))



def compute_initial_confidence(emotion_data, speech_data):
    initial = (
        0.35 * emotion_data["facial_confidence"]
        + 0.30 * speech_data["vocal_confidence"]
        + 0.20 * speech_data["fluency_score"]
        + 0.15 * speech_data["content_delivery"]
    )
    return clamp(initial)



def compute_final_confidence(initial_score, answer_scores):
    if answer_scores:
        answer_quality = float(np.mean(answer_scores))
    else:
        answer_quality = 0.0

    improvement_bonus = clamp((answer_quality - 50.0) * 0.6 + 20.0)
    final = (0.60 * initial_score) + (0.25 * answer_quality) + (0.15 * improvement_bonus)
    final = max(final, initial_score + 2.0 if answer_quality > 55 else final)
    return clamp(final), answer_quality, improvement_bonus


# =========================================================
# FEEDBACK GENERATION
# =========================================================
def generate_feedback(weak_areas, speech_data, emotion_data, answer_evals):
    strengths = []
    improvements = []

    if speech_data.get("duration_seconds", 0) > 1200:
        improvements.append("Your presentation exceeded the 20-minute limit. Consider skipping less critical slides, streamlining your transitions, and reducing filler words to keep it concise.")

    if emotion_data["facial_confidence"] >= 60:
        strengths.append("Facial delivery looks reasonably stable across the presentation.")
    if speech_data["fluency_score"] >= 60:
        strengths.append("Speech fluency is acceptable, with manageable hesitation.")
    if speech_data["content_delivery"] >= 55:
        strengths.append("Your spoken explanation covers a fair amount of the slide content.")

    if "nervousness" in weak_areas:
        improvements.append("Work on steadier visual delivery and reduce nervous facial spikes during important explanations.")
    if "fluency" in weak_areas:
        improvements.append("Reduce filler words and maintain a more consistent speaking pace.")
    if "content_depth" in weak_areas:
        improvements.append("Explain implementation details with more depth, not just high-level outcomes.")
    if "technical_understanding" in weak_areas:
        improvements.append("Prepare stronger justifications for design choices, trade-offs, and failure cases.")
    if "retrieval_gap" in weak_areas:
        improvements.append("Align your spoken explanation more closely with the core concepts shown on your slides.")

    avg_answer_score = float(np.mean([x["score"] for x in answer_evals])) if answer_evals else 0.0
    if avg_answer_score >= 70:
        strengths.append("Your viva answers show good recovery and stronger confidence after feedback.")
    elif avg_answer_score > 0:
        improvements.append("During viva responses, try to be more precise and connect every answer back to your project pipeline.")

    if not strengths:
        strengths.append("You have a workable presentation base and can improve quickly with targeted viva practice.")
    if not improvements:
        improvements.append("Keep practicing concise technical explanations with clearer module-level justification.")

    return strengths, improvements


# =========================================================
# UI HELPERS
# =========================================================
def render_metric_card(label, value, suffix=""):
    st.metric(label, f"{value:.2f}{suffix}" if isinstance(value, (float, int)) else str(value))


def render_context(context_list):
    for item in context_list:
        with st.expander(f"{item['source']} · relevance {item['score']:.3f}"):
            st.write(item["text"])


# =========================================================
# MAIN APP
# =========================================================
def main():
    st.title("GenAI Viva Coach")
    st.caption("Upload a capstone PPT/PDF and a presentation video. The app computes an initial confidence score, runs a retrieval-grounded viva, and shows post-viva improvement.")

    if "phase" not in st.session_state:
        st.session_state.phase = 1
        st.session_state.results_1 = None
        st.session_state.questions_1 = []
        st.session_state.evals_1 = {}
        st.session_state.results_2 = None
        st.session_state.questions_2 = []
        st.session_state.evals_2 = {}
        st.session_state.slide_text_full = ""
        st.session_state.slide_chunks = []

    # UI Flow based on Phase
    if st.session_state.phase == 1:
        st.header("Phase 1: Initial Upload")
        slides_file = st.file_uploader("Upload capstone slides (.pptx or .pdf)", type=["pptx", "pdf"], key="slides")
        video_file = st.file_uploader("Upload presentation video (.mp4, .mov, .avi)", type=["mp4", "mov", "avi", "mkv"], key="video_1")
        if st.button("Run Initial Analysis"):
            if slides_file is None or video_file is None:
                st.error("Please upload both the slides file and the presentation video.")
                return
            
            with st.spinner("Processing files..."):
                slides_suffix = ".pptx" if slides_file.name.lower().endswith(".pptx") else ".pdf"
                slides_path = save_uploaded_file(slides_file, slides_suffix)
                video_suffix = os.path.splitext(video_file.name)[1] or ".mp4"
                video_path = save_uploaded_file(video_file, video_suffix)
                audio_path = tempfile.NamedTemporaryFile(delete=False, suffix=".wav").name

                if slides_suffix == ".pptx":
                    slide_chunks = extract_text_from_pptx(slides_path)
                else:
                    slide_chunks = extract_text_from_pdf(slides_path)
                st.session_state.slide_chunks = slide_chunks
                st.session_state.slide_text_full = "\n\n".join([c["text"] for c in slide_chunks])

                extract_audio_from_video(video_path, audio_path)
                transcript_text, segments = transcribe_audio(audio_path)
                speech_data = analyze_speech(audio_path, transcript_text, segments, st.session_state.slide_text_full)
                emotion_data = analyze_video_emotions(video_path)

                transcript_chunks = chunk_transcript(segments)
                all_chunks = slide_chunks + transcript_chunks
                if not all_chunks:
                    all_chunks = [{"source": "Fallback", "text": transcript_text or st.session_state.slide_text_full or "No content extracted."}]
                embedder, index, _ = build_vector_index(all_chunks)

                weak_areas = detect_weak_areas(emotion_data, speech_data)
                initial_score = compute_initial_confidence(emotion_data, speech_data)
                questions = generate_viva_questions(st.session_state.slide_text_full + "\n" + transcript_text, weak_areas, embedder, index, all_chunks, api_key=GEMINI_API_KEY_1)

                st.session_state.questions_1 = questions
                st.session_state.results_1 = {
                    "transcript_text": transcript_text,
                    "segments": segments,
                    "speech_data": speech_data,
                    "emotion_data": emotion_data,
                    "weak_areas": weak_areas,
                    "initial_score": initial_score,
                    "embedder": embedder,
                    "index": index,
                    "all_chunks": all_chunks,
                }
                st.session_state.phase = 2
                st.rerun()

    elif st.session_state.phase == 2:
        st.header("Phase 2: Initial Viva Simulation")
        st.write("Your video has been analyzed. Please answer the following viva questions to complete your initial evaluation.")
        
        questions = st.session_state.questions_1
        embedder = st.session_state.results_1["embedder"]
        answers = []
        for i, q in enumerate(questions, start=1):
            st.markdown(f"**Q{i}. [{q['type']}] {q['question']}**")
            ans = st.text_area(f"Your answer for Q{i}", key=f"ans_1_{i}", height=100)
            answers.append(ans)

        if st.button("Evaluate Viva Answers"):
            answer_evals = []
            answer_scores = []
            for q, ans in zip(questions, answers):
                evaluation = evaluate_answer(
                    answer=ans, 
                    question_obj=q, 
                    embedder=st.session_state.results_1["embedder"], 
                    index=st.session_state.results_1["index"], 
                    chunks=st.session_state.results_1["all_chunks"],
                    slide_text=st.session_state.slide_text_full,
                    api_key=GEMINI_API_KEY_1
                )
                answer_evals.append(evaluation)
                answer_scores.append(evaluation["score"])
            
            initial_score_raw = st.session_state.results_1["initial_score"]
            orig_conf_score, answer_quality, improvement_bonus = compute_final_confidence(initial_score_raw, answer_scores)
            
            st.session_state.evals_1 = {
                "answer_evals": answer_evals,
                "original_confidence_score": orig_conf_score,
                "answer_quality": answer_quality
            }
            st.session_state.phase = 3
            st.rerun()

    elif st.session_state.phase == 3:
        st.header("Phase 3: Feedback & Second Attempt")
        res1 = st.session_state.results_1
        evals1 = st.session_state.evals_1
        
        orig_score = evals1["original_confidence_score"]
        st.metric("Original Total Confidence Score", f"{orig_score:.2f}")

        st.subheader("Initial Metrics")
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Fluency Score", f"{res1['speech_data']['fluency_score']:.2f}")
        c2.metric("Vocal Confidence", f"{res1['speech_data']['vocal_confidence']:.2f}")
        c3.metric("Facial Confidence", f"{res1['emotion_data']['facial_confidence']:.2f}")
        c4.metric("Words / Minute", f"{res1['speech_data']['words_per_minute']:.2f}")
        
        strengths, improvements = generate_feedback(res1["weak_areas"], res1["speech_data"], res1["emotion_data"], evals1["answer_evals"])
        
        st.markdown("### Strengths")
        for s in strengths:
            st.write(f"- {s}")
        st.markdown("### Areas to improve")
        for imp in improvements:
            st.write(f"- {imp}")
            
        st.markdown("### Detailed Viva Feedback")
        for i, eval_obj in enumerate(evals1["answer_evals"], start=1):
            with st.expander(f"Question {i} Feedback (Score: {eval_obj['score']:.1f})"):
                st.write(eval_obj["feedback"])

        st.divider()
        st.write("You can now record an improved presentation based on the feedback above and upload the new video.")
        video_file_2 = st.file_uploader("Upload Improved Video (.mp4, .mov, .avi)", type=["mp4", "mov", "avi", "mkv"], key="video_2")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Skip and Finish"):
                st.session_state.phase = 5
                st.rerun()
        with col2:
            if st.button("Run New Analysis"):
                if video_file_2 is None:
                    st.error("Please upload the improved video.")
                else:
                    with st.spinner("Processing new video..."):
                        video_suffix = os.path.splitext(video_file_2.name)[1] or ".mp4"
                        video_path = save_uploaded_file(video_file_2, video_suffix)
                        audio_path = tempfile.NamedTemporaryFile(delete=False, suffix=".wav").name
                        
                        extract_audio_from_video(video_path, audio_path)
                        transcript_text, segments = transcribe_audio(audio_path)
                        speech_data = analyze_speech(audio_path, transcript_text, segments, st.session_state.slide_text_full)
                        emotion_data = analyze_video_emotions(video_path)

                        transcript_chunks = chunk_transcript(segments)
                        all_chunks = st.session_state.slide_chunks + transcript_chunks
                        if not all_chunks:
                            all_chunks = [{"source": "Fallback", "text": transcript_text or st.session_state.slide_text_full or "No content extracted."}]
                        embedder, index, _ = build_vector_index(all_chunks)

                        weak_areas = detect_weak_areas(emotion_data, speech_data)
                        initial_score = compute_initial_confidence(emotion_data, speech_data)
                        questions = generate_viva_questions(st.session_state.slide_text_full + "\n" + transcript_text, weak_areas, embedder, index, all_chunks, api_key=GEMINI_API_KEY_2)
                        
                        st.session_state.questions_2 = questions
                        st.session_state.results_2 = {
                            "transcript_text": transcript_text,
                            "segments": segments,
                            "speech_data": speech_data,
                            "emotion_data": emotion_data,
                            "weak_areas": weak_areas,
                            "initial_score": initial_score,
                            "embedder": embedder,
                            "index": index,
                            "all_chunks": all_chunks,
                        }
                        st.session_state.phase = 4
                        st.rerun()
                    
    elif st.session_state.phase == 4:
        st.header("Phase 4: Second Viva Simulation")
        st.write("Your improved video has been analyzed. A new set of questions has been generated to verify your progress.")
        
        questions = st.session_state.questions_2
        embedder = st.session_state.results_2["embedder"]
        answers = []
        for i, q in enumerate(questions, start=1):
            st.markdown(f"**Q{i}. [{q['type']}] {q['question']}**")
            ans = st.text_area(f"Your answer for Q{i}", key=f"ans_2_{i}", height=100)
            answers.append(ans)

        if st.button("Evaluate New Viva Answers"):
            answer_evals = []
            answer_scores = []
            for q, ans in zip(questions, answers):
                evaluation = evaluate_answer(
                    answer=ans, 
                    question_obj=q, 
                    embedder=st.session_state.results_2["embedder"], 
                    index=st.session_state.results_2["index"], 
                    chunks=st.session_state.results_2["all_chunks"],
                    slide_text=st.session_state.slide_text_full,
                    api_key=GEMINI_API_KEY_2
                )
                answer_evals.append(evaluation)
                answer_scores.append(evaluation["score"])
            
            initial_score_raw = st.session_state.results_2["initial_score"]
            new_conf_score, answer_quality, improvement_bonus = compute_final_confidence(initial_score_raw, answer_scores)
            
            st.session_state.evals_2 = {
                "answer_evals": answer_evals,
                "new_confidence_score": new_conf_score,
                "answer_quality": answer_quality
            }
            st.session_state.phase = 5
            st.rerun()

    elif st.session_state.phase == 5:
        st.header("Phase 5: Final Evaluation Overview")
        if st.session_state.results_2 is None:
            st.write("You skipped the second video upload. Here is your summary from the first round:")
            st.metric("Final Confidence Score", f"{st.session_state.evals_1['original_confidence_score']:.2f}")
        else:
            orig_score = st.session_state.evals_1["original_confidence_score"]
            new_score = st.session_state.evals_2["new_confidence_score"]
            delta = new_score - orig_score
            
            st.subheader("Score Comparison")
            c1, c2, c3 = st.columns(3)
            c1.metric("Original Score", f"{orig_score:.2f}")
            c2.metric("New Score", f"{new_score:.2f}", delta=f"{delta:.2f}")
            
            if delta > 0:
                st.success("Great job! Your confidence score improved.")
            else:
                st.warning("Your score didn't improve. Keep practicing!")
                
            st.subheader("New Metrics")
            res2 = st.session_state.results_2
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Fluency Score", f"{res2['speech_data']['fluency_score']:.2f}")
            col2.metric("Vocal Confidence", f"{res2['speech_data']['vocal_confidence']:.2f}")
            col3.metric("Facial Confidence", f"{res2['emotion_data']['facial_confidence']:.2f}")
            col4.metric("Words / Minute", f"{res2['speech_data']['words_per_minute']:.2f}")
            
            strengths, improvements = generate_feedback(res2["weak_areas"], res2["speech_data"], res2["emotion_data"], st.session_state.evals_2["answer_evals"])
            st.markdown("### Final Strengths")
            for s in strengths:
                st.write(f"- {s}")
            st.markdown("### Final Areas to improve")
            for imp in improvements:
                st.write(f"- {imp}")

            st.markdown("### Detailed Viva Feedback")
            for i, eval_obj in enumerate(st.session_state.evals_2["answer_evals"], start=1):
                with st.expander(f"Question {i} Feedback (Score: {eval_obj['score']:.1f})"):
                    st.write(eval_obj["feedback"])
        
        if st.button("Start Over"):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()


if __name__ == "__main__":
    main()
